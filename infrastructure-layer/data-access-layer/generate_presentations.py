#!/usr/bin/env python3
"""
Script to generate PowerPoint and PDF presentations from README.md
"""

import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Preformatted
from reportlab.lib.enums import TA_LEFT, TA_CENTER
from reportlab.lib.colors import HexColor

def parse_markdown(filename):
    """Parse markdown file and extract sections"""
    with open(filename, 'r', encoding='utf-8') as f:
        content = f.read()
    
    sections = []
    current_section = None
    current_content = []
    
    lines = content.split('\n')
    i = 0
    
    while i < len(lines):
        line = lines[i]
        
        # Main headers (##)
        if line.startswith('## '):
            if current_section:
                sections.append({
                    'title': current_section,
                    'content': '\n'.join(current_content).strip(),
                    'level': 1
                })
            current_section = line[3:].strip()
            current_content = []
        
        # Sub headers (###)
        elif line.startswith('### '):
            if current_content:
                sections.append({
                    'title': current_section,
                    'content': '\n'.join(current_content).strip(),
                    'level': 1
                })
                current_content = []
            
            subsection = line[4:].strip()
            sections.append({
                'title': subsection,
                'content': '',
                'level': 2,
                'parent': current_section
            })
        
        # Sub-sub headers (####)
        elif line.startswith('#### '):
            subsubsection = line[5:].strip()
            current_content.append(f"\n{subsubsection}\n")
        
        # Code blocks
        elif line.startswith('```'):
            code_lines = []
            i += 1
            while i < len(lines) and not lines[i].startswith('```'):
                code_lines.append(lines[i])
                i += 1
            current_content.append('\n[CODE]\n' + '\n'.join(code_lines[:15]) + '\n[/CODE]\n')
        
        # Bullet points
        elif line.startswith('- ') or line.startswith('* '):
            current_content.append(line)
        
        # Regular text
        elif line.strip() and not line.startswith('---'):
            current_content.append(line)
        
        i += 1
    
    # Add last section
    if current_section:
        sections.append({
            'title': current_section,
            'content': '\n'.join(current_content).strip(),
            'level': 1
        })
    
    return sections

def create_powerpoint(sections, output_file):
    """Create PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Data Access Layer dengan Dapper"
    subtitle.text = "Panduan Lengkap Implementasi Dapper di .NET"
    
    # Add content slides
    for section in sections:
        if section['level'] == 1:
            # Section title slide
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = section['title']
            
            # Process content
            text_content = section['content']
            
            # Handle code blocks
            if '[CODE]' in text_content:
                code_match = re.search(r'\[CODE\](.*?)\[/CODE\]', text_content, re.DOTALL)
                if code_match:
                    code = code_match.group(1).strip()
                    # Add code box
                    tf = content.text_frame
                    tf.clear()
                    p = tf.paragraphs[0]
                    p.text = code[:500]  # Limit code length
                    p.font.size = Pt(10)
                    p.font.name = 'Courier New'
            else:
                # Regular content
                tf = content.text_frame
                tf.clear()
                
                # Split content into bullet points or paragraphs
                lines = [l for l in text_content.split('\n') if l.strip()]
                for line in lines[:10]:  # Limit to 10 lines per slide
                    p = tf.add_paragraph()
                    
                    # Clean markdown formatting
                    clean_line = line.replace('**', '').replace('`', '').strip()
                    
                    # Check if it's a bullet point
                    if clean_line.startswith('- ') or clean_line.startswith('* '):
                        p.text = clean_line[2:]
                        p.level = 0
                    elif clean_line.startswith('  - ') or clean_line.startswith('  * '):
                        p.text = clean_line[4:]
                        p.level = 1
                    else:
                        p.text = clean_line
                    
                    p.font.size = Pt(14)
    
    prs.save(output_file)
    print(f"PowerPoint presentation saved to: {output_file}")

def create_pdf(sections, output_file):
    """Create PDF presentation"""
    doc = SimpleDocTemplate(output_file, pagesize=letter)
    story = []
    styles = getSampleStyleSheet()
    
    # Custom styles
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        textColor=HexColor('#2E4057'),
        spaceAfter=30,
        alignment=TA_CENTER
    )
    
    heading_style = ParagraphStyle(
        'CustomHeading',
        parent=styles['Heading2'],
        fontSize=18,
        textColor=HexColor('#2E4057'),
        spaceAfter=12,
        spaceBefore=12
    )
    
    subheading_style = ParagraphStyle(
        'CustomSubHeading',
        parent=styles['Heading3'],
        fontSize=14,
        textColor=HexColor('#4A5F7F'),
        spaceAfter=10,
        spaceBefore=10
    )
    
    body_style = ParagraphStyle(
        'CustomBody',
        parent=styles['BodyText'],
        fontSize=11,
        spaceAfter=8
    )
    
    code_style = ParagraphStyle(
        'Code',
        parent=styles['Code'],
        fontSize=9,
        leftIndent=20,
        fontName='Courier'
    )
    
    # Title page
    story.append(Spacer(1, 2*inch))
    story.append(Paragraph("Data Access Layer dengan Dapper", title_style))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("Panduan Lengkap Implementasi Dapper di .NET", styles['Normal']))
    story.append(PageBreak())
    
    # Content
    for section in sections:
        if section['level'] == 1:
            story.append(Paragraph(section['title'], heading_style))
            story.append(Spacer(1, 0.2*inch))
            
            content = section['content']
            
            # Handle code blocks
            if '[CODE]' in content:
                # Extract code
                parts = re.split(r'\[CODE\](.*?)\[/CODE\]', content, flags=re.DOTALL)
                for i, part in enumerate(parts):
                    if i % 2 == 0:  # Regular text
                        if part.strip():
                            lines = [l for l in part.split('\n') if l.strip()]
                            for line in lines:
                                # Remove all markdown formatting
                                clean_line = line.replace('**', '').replace('`', '').strip()
                                # Escape special characters
                                clean_line = clean_line.replace('<', '&lt;').replace('>', '&gt;')
                                clean_line = clean_line.replace('&', '&amp;')
                                
                                if clean_line.startswith('- ') or clean_line.startswith('* '):
                                    story.append(Paragraph('• ' + clean_line[2:], body_style))
                                elif clean_line:
                                    story.append(Paragraph(clean_line, body_style))
                    else:  # Code block
                        code_lines = part.strip().split('\n')[:20]  # Limit code lines
                        for code_line in code_lines:
                            story.append(Preformatted(code_line, code_style))
                        story.append(Spacer(1, 0.1*inch))
            else:
                # Regular content
                lines = [l for l in content.split('\n') if l.strip()]
                for line in lines:
                    # Remove all markdown formatting
                    clean_line = line.replace('**', '').replace('`', '').strip()
                    # Escape special characters
                    clean_line = clean_line.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                    
                    if clean_line.startswith('- ') or clean_line.startswith('* '):
                        story.append(Paragraph('• ' + clean_line[2:], body_style))
                    elif clean_line.startswith('  - ') or clean_line.startswith('  * '):
                        story.append(Paragraph('  • ' + clean_line[4:], body_style))
                    elif clean_line:
                        story.append(Paragraph(clean_line, body_style))
            
            story.append(Spacer(1, 0.3*inch))
        
        elif section['level'] == 2:
            story.append(Paragraph(section['title'], subheading_style))
            story.append(Spacer(1, 0.1*inch))
    
    doc.build(story)
    print(f"PDF presentation saved to: {output_file}")

def main():
    """Main function"""
    readme_file = 'README.md'
    pptx_output = 'Data-Access-Layer-Presentation.pptx'
    pdf_output = 'Data-Access-Layer-Presentation.pdf'
    
    print("Parsing README.md...")
    sections = parse_markdown(readme_file)
    print(f"Found {len(sections)} sections")
    
    print("\nGenerating PowerPoint presentation...")
    create_powerpoint(sections, pptx_output)
    
    print("\nGenerating PDF presentation...")
    create_pdf(sections, pdf_output)
    
    print("\n✅ Presentations generated successfully!")
    print(f"   - PowerPoint: {pptx_output}")
    print(f"   - PDF: {pdf_output}")

if __name__ == '__main__':
    main()
